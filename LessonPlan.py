import streamlit as st
from ai import ask_ai
from docx import Document
from io import BytesIO
from pptx import Presentation


# ---------------- CLEAN DISPLAY ----------------
def clean_lesson_plan_display(text):

    lines = text.split("\n")
    cleaned = []

    for line in lines:
        if "-------------------------------------" in line:
            cleaned.append("")
        else:
            cleaned.append(line)

    return "\n".join(cleaned)


# ---------------- GENERATE LESSON PLAN ----------------
def generate_lesson_plan(country, class_level, subject, topic, duration, sessions):

    system_prompt = """
You are an expert teacher trainer familiar with the Ugandan education system
and the Competency-Based Curriculum (CBC).

Generate professional lesson plans similar to those used in Ugandan schools.
"""

    user_prompt = f"""
Country: {country}
Curriculum: Uganda CBC
Class Level: {class_level}
Subject: {subject}

Topic:
{topic}

Session Duration: {duration} hours
Number of Sessions: {sessions}

Create a structured professional lesson plan.

Use this structure:

-------------------------------------
LESSON IDENTIFICATION
School:
Teacher:
Subject:
Class:
Duration:
Topic:
Sub-topic:

-------------------------------------
COMPETENCY
Key competency to be developed

-------------------------------------
LEARNING OBJECTIVES
• objective
• objective

-------------------------------------
KEY INQUIRY QUESTION

-------------------------------------
LEARNING ACTIVITIES
Teacher Activities
Learner Activities

-------------------------------------
ASSESSMENT METHODS

-------------------------------------
LEARNING RESOURCES

-------------------------------------
TEACHER REFLECTION

-------------------------------------

Ensure explanations match the cognitive level of {class_level} learners in Uganda.
"""

    return ask_ai(system_prompt, user_prompt)


# ---------------- GENERATE SLIDE CONTENT ----------------
def generate_slide_content(country, class_level, subject, topic):

    system_prompt = """
You are a professional teacher and presentation designer.
Generate teaching slides appropriate for classroom instruction.
"""

    user_prompt = f"""
Country: {country}
Class Level: {class_level}
Subject: {subject}

Topic:
{topic}

Create structured teaching slides.

Rules:
• Maximum 5 bullet points
• Clear titles
• Professional explanations
• Include speaker notes

Use the format:

Slide 1
Title: ...
Content:
- point
- point
Notes:
teacher explanation

Slide 2
Title: ...
Content:
- point
- point
Notes:
teacher explanation
"""

    return ask_ai(system_prompt, user_prompt)


# ---------------- CREATE DOCX ----------------
def generate_docx(text):

    doc = Document()
    doc.add_heading("Lesson Plan", level=0)

    lines = text.split("\n")

    for line in lines:

        if "-------------------------------------" in line:
            doc.add_paragraph("")
        else:
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    return buffer


# ---------------- CREATE PPT ----------------
def generate_ppt(slide_text, topic):

    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Teaching Slides"

    slides = slide_text.split("Slide")

    for block in slides:

        if "Title:" not in block:
            continue

        lines = block.split("\n")

        title = ""
        content = []
        notes = ""
        mode = None

        for line in lines:

            if line.startswith("Title:"):
                title = line.replace("Title:", "").strip()

            elif line.startswith("Content"):
                mode = "content"

            elif line.startswith("Notes"):
                mode = "notes"

            elif line.startswith("-") and mode == "content":
                content.append(line.replace("-", "").strip())

            elif mode == "notes":
                notes += line + " "

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)

        slide_obj.shapes.title.text = title

        tf = slide_obj.placeholders[1].text_frame
        tf.clear()

        for point in content:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

        slide_obj.notes_slide.notes_text_frame.text = notes

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer


# ---------------- STREAMLIT UI ----------------
def lessonplan():

    st.subheader("AI Lesson Plan Generator")

    country = "Uganda"

    st.divider()

    education_level = st.selectbox(
        "Education Level",
        ["Primary", "Secondary"]
    )

    if education_level == "Primary":

        class_level = st.selectbox(
            "Class",
            ["Primary 1","Primary 2","Primary 3","Primary 4","Primary 5","Primary 6","Primary 7"]
        )

        subject = st.selectbox(
            "Subject",
            [
                "English","Mathematics","Science","Social Studies",
                "Christian Religious Education",
                "Islamic Religious Education",
                "Creative Arts & Physical Education"
            ]
        )

    else:

        class_level = st.selectbox(
            "Class",
            ["Senior 1","Senior 2","Senior 3","Senior 4","Senior 5","Senior 6"]
        )

        if class_level in ["Senior 1","Senior 2","Senior 3","Senior 4"]:

            subject = st.selectbox(
                "Subject",
                [
                    "English","Mathematics","Biology","Chemistry","Physics",
                    "Geography","History","ICT","Entrepreneurship",
                    "Agriculture","Fine Art","Technical Drawing",
                    "Nutrition and Technology","French","Luganda"
                ]
            )

        else:

            subject = st.selectbox(
                "Subject",
                [
                    "Physics","Chemistry","Biology","Mathematics",
                    "Agriculture","Food and Nutrition","Technical Drawing",
                    "History","Geography","Economics",
                    "Literature in English","Divinity","Entrepreneurship",
                    "Fine Art","Luganda","French","German","Arabic",
                    "General Paper","Subsidiary Mathematics","Subsidiary ICT"
                ]
            )

    topic = st.text_area("Topic / Unit Details", height=150)

    duration = st.number_input("Session Duration (hours)", min_value=1)
    sessions = st.number_input("Number of Sessions", min_value=1)

    st.divider()

    col1, col2 = st.columns(2)

    # -------- LESSON PLAN --------
    with col1:

        st.markdown("### Lesson Plan")
        st.markdown("---")

        if st.button("Generate Lesson Plan"):

            with st.spinner("Generating lesson plan..."):

                lesson_plan = generate_lesson_plan(
                    country,
                    class_level,
                    subject,
                    topic,
                    duration,
                    sessions
                )

            st.session_state.lesson_plan = lesson_plan

        if "lesson_plan" in st.session_state:

            cleaned_plan = clean_lesson_plan_display(st.session_state.lesson_plan)

            st.text(cleaned_plan)

            docx = generate_docx(st.session_state.lesson_plan)

            st.download_button(
                "Download Lesson Plan (DOCX)",
                data=docx,
                file_name="lesson_plan.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

    # -------- SLIDES --------
    with col2:

        st.markdown("### Teaching Slides")
        st.markdown("---")

        if st.button("Generate Teaching Slides"):

            with st.spinner("Generating slides..."):

                slides = generate_slide_content(
                    country,
                    class_level,
                    subject,
                    topic
                )

            st.session_state.slides = slides

        if "slides" in st.session_state:

            st.markdown(st.session_state.slides)

            ppt = generate_ppt(st.session_state.slides, topic)

            st.download_button(
                "Download Slides (PPTX)",
                data=ppt,
                file_name="lesson_slides.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )